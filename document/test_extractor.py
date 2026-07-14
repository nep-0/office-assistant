import unittest
import zipfile
from io import BytesIO
from unittest.mock import patch

from openpyxl import Workbook
from pptx import Presentation

from extractor import ExtractionError, call_ocr, extract_upload, ocr_error_message, ocr_header_filename


class ExtractorTests(unittest.TestCase):
    def test_extracts_docx_text(self):
        package = extract_upload("report.docx", make_zip({
            "word/document.xml": """
                <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
                  <w:body><w:p><w:r><w:t>Hello DOCX</w:t></w:r></w:p></w:body>
                </w:document>
            """,
        }))

        self.assertEqual(package["schema_version"], "v1")
        self.assertIn("Hello DOCX", package["markdown"])
        self.assertEqual(package["source_anchors"][0]["kind"], "section")

    def test_extracts_xlsx_rows(self):
        workbook = Workbook()
        sheet = workbook.active
        sheet["A1"] = "Name"
        sheet["B1"] = "Total"
        buffer = BytesIO()
        workbook.save(buffer)

        package = extract_upload("table.xlsx", buffer.getvalue())

        self.assertIn("Name", package["markdown"])
        self.assertIn("Total", package["markdown"])
        self.assertEqual(package["source_anchors"][0]["kind"], "sheet")

    def test_extracts_pptx_slide_text(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Hello PPTX"
        buffer = BytesIO()
        presentation.save(buffer)

        package = extract_upload("slides.pptx", buffer.getvalue())

        self.assertIn("Hello PPTX", package["markdown"])
        self.assertEqual(package["source_anchors"][0]["kind"], "slide")

    def test_extracts_pdf_native_text(self):
        package = extract_upload("file.pdf", make_pdf("Hello PDF native text for extraction"))

        self.assertIn("Hello PDF native text", package["markdown"])
        self.assertEqual(package["source_anchors"][0]["kind"], "page")

    def test_image_uses_ocr_fallback(self):
        package = extract_upload(
            "scan.png",
            b"fake-image",
            ocr_func=lambda _data, filename: f"OCR text from {filename}",
        )

        self.assertIn("OCR text from scan.png", package["markdown"])
        self.assertTrue(package["ocr"]["used"])
        self.assertEqual(package["source_anchors"][0]["kind"], "image")

    def test_image_preserves_chinese_filename_in_metadata(self):
        package = extract_upload(
            "扫描.png",
            b"fake-image",
            ocr_func=lambda _data, filename: f"OCR text from {filename}",
        )

        self.assertEqual(package["metadata"]["filename"], "扫描.png")
        self.assertEqual(package["source_anchors"][0]["label"], "扫描.png")

    def test_ocr_header_filename_falls_back_for_chinese_names(self):
        self.assertEqual(ocr_header_filename("扫描.png"), "upload.png")
        self.assertEqual(ocr_header_filename("scan.png"), "scan.png")

    def test_unsupported_file_fails_clearly(self):
        with self.assertRaises(ExtractionError) as caught:
            extract_upload("notes.txt", b"hello")

        self.assertEqual(caught.exception.code, "unsupported_office_input")

    def test_ocr_error_message_uses_service_payload(self):
        message = ocr_error_message('{"code":"ppocr_failed","message":"cache is read-only"}', "Bad Gateway")

        self.assertEqual(message, "cache is read-only")

    def test_ocr_timeout_is_structured_extraction_error(self):
        with patch("urllib.request.urlopen", side_effect=TimeoutError()):
            with self.assertRaises(ExtractionError) as caught:
                call_ocr(b"image", "scan.png", "http://ocr.test", None)

        self.assertEqual(caught.exception.code, "ocr_timeout")


def make_zip(files):
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        for name, content in files.items():
            archive.writestr(name, content)
    return buffer.getvalue()


def make_pdf(text):
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> /MediaBox [0 0 612 792] /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("latin-1")
    objects.append(b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream")

    output = BytesIO()
    output.write(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(output.tell())
        output.write(f"{index} 0 obj\n".encode("ascii"))
        output.write(obj)
        output.write(b"\nendobj\n")
    xref = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.write(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.write(b"trailer\n")
    output.write(f"<< /Size {len(objects) + 1} /Root 1 0 R >>\n".encode("ascii"))
    output.write(b"startxref\n")
    output.write(str(xref).encode("ascii"))
    output.write(b"\n%%EOF\n")
    return output.getvalue()


if __name__ == "__main__":
    unittest.main()
